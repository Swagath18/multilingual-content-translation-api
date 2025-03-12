# backend/utils/text_extraction.py
from fastapi import UploadFile, HTTPException
from pptx import Presentation
from io import BytesIO
from bs4 import BeautifulSoup
import PyPDF2
import logging

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

async def extract_text_from_pdf(pdf_file: UploadFile) -> str:
    """Extract text from a PDF file asynchronously."""
    try:
        logger.info(f"Extracting text from PDF: {pdf_file.filename}")
        pdf_reader = PyPDF2.PdfReader(pdf_file.file)
        text = ""
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error reading PDF: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Error reading PDF: {str(e)}")
    finally:
        pdf_file.file.close()

async def extract_text_from_ppt(ppt_file: UploadFile) -> str:
    """Extract text from a PPT file asynchronously."""
    try:
        logger.info(f"Extracting text from PPT: {ppt_file.filename}")
        ppt_file_content = await ppt_file.read()
        ppt_file_buffer = BytesIO(ppt_file_content)
        ppt = Presentation(ppt_file_buffer)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error reading PPT: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Error reading PPT: {str(e)}")
    finally:
        ppt_file.file.close()

async def extract_text_from_html(html_file: UploadFile) -> str:
    """Extract text from an HTML file asynchronously."""
    try:
        logger.info(f"Extracting text from HTML: {html_file.filename}")
        html_content = html_file.file.read().decode("utf-8")
        soup = BeautifulSoup(html_content, "html.parser")
        text = soup.get_text(separator="\n")
        return text
    except Exception as e:
        logger.error(f"Error reading HTML: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Error reading HTML: {str(e)}")
    finally:
        html_file.file.close()

async def extract_text_from_txt(txt_file: UploadFile) -> str:
    """Extract text from a TXT file asynchronously."""
    try:
        logger.info(f"Extracting text from TXT: {txt_file.filename}")
        file_content = await txt_file.read()
        text = file_content.decode('utf-8')
        return text
    except Exception as e:
        logger.error(f"Error reading TXT file: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Error reading TXT file: {str(e)}")
    finally:
        txt_file.file.close()